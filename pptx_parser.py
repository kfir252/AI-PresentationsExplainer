from pptx import Presentation

class PresentationParser:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.presentation_data = None
        self.total_slides = 0  
        
        presentation = Presentation(self.pptx_path)
        self.total_slides = len(presentation.slides)
        self.presentation_data = [
            {
                "slide_number": slide_number,
                "title": title if (title := next((shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape == slide.shapes.title), "")) else "",
                "content": [shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape != slide.shapes.title]
            }
            for slide_number, slide in enumerate(presentation.slides, start=1)
        ]

    
    def get_slide(self, slide_index:int ) -> dict[str, None]:
        return next(( slide for slide in self.presentation_data if slide['slide_number'] == slide_index ),None,)
